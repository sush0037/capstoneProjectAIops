"""Embed real app screenshots into ClaimIQ_Presentation.pptx (Slide 9)
and ClaimIQ_Project_Report.docx (Section 11)."""

import pathlib, copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Inches as DInches, Pt as DPt, RGBColor as DRGBColor
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

BASE   = pathlib.Path(__file__).parent.parent
SS_DIR = BASE / "data" / "screenshots"
PPTX   = BASE / "ClaimIQ_Presentation.pptx"
DOCX   = BASE / "ClaimIQ_Project_Report.docx"

# ── The 5 best screenshots with captions ───────────────────────────────────
SCREENSHOTS = [
    ("01_home.png",              "1. Home — Upload & System Status"),
    ("03_pipeline_complete.png", "2. AI Pipeline — OCR Extraction Results"),
    ("06_ai_recommendation.png", "3. AI Recommendation & Policy Citations"),
    ("08_claims_history.png",    "4. Claims History Dashboard"),
    ("09_about.png",             "5. About & System Information"),
]

BLUE_DARK  = RGBColor(0x0A, 0x6E, 0xBD)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLUE_LIGHT = RGBColor(0xE8, 0xF4, 0xFD)
DARK       = RGBColor(0x1A, 0x1A, 0x2E)
GREY       = RGBColor(0x55, 0x55, 0x66)


# ════════════════════════════════════════════════════════════════════════════
# 1.  UPDATE PPTX — Slide 9 (index 8)
# ════════════════════════════════════════════════════════════════════════════
print("Updating PPTX Slide 9 (Demo Screenshots)...")

prs = Presentation(str(PPTX))
slide = prs.slides[8]   # 0-based index, slide 9

# Remove all shapes except the header bar and its text (first 3 shapes)
# Keep shapes whose top is < 1.3" (the header region)
to_remove = []
for shape in slide.shapes:
    if shape.top > Inches(1.3):
        to_remove.append(shape)
for shape in to_remove:
    sp = shape._element
    sp.getparent().remove(sp)

# ── Layout: 3 images top row, 2 images bottom row ──────────────────────────
IMG_W  = Inches(4.0)
IMG_H  = Inches(2.5)
GAP_X  = Inches(0.155)
CAP_H  = Inches(0.3)
ROW1_Y = Inches(1.32)
ROW2_Y = Inches(4.3)

# Row 1 x positions (3 images)
row1_xs = [Inches(0.3), Inches(4.455), Inches(8.61)]
# Row 2 x positions (2 images, centered): total = 2*4.0 + gap = 8.155", offset = (13.33-8.155)/2
row2_start = Inches((13.33 - 8.155) / 2)
row2_xs = [row2_start, row2_start + IMG_W + GAP_X]

all_positions = [(x, ROW1_Y) for x in row1_xs] + [(x, ROW2_Y) for x in row2_xs]

for i, ((filename, caption), (x, y)) in enumerate(zip(SCREENSHOTS, all_positions)):
    img_path = SS_DIR / filename
    if not img_path.exists():
        print(f"  MISSING: {filename}")
        continue

    # Add screenshot image
    pic = slide.shapes.add_picture(str(img_path), x, y, IMG_W, IMG_H)

    # White border around image
    from pptx.oxml.ns import qn as pqn
    from lxml import etree
    spPr = pic._element.spPr
    ln = etree.SubElement(spPr, pqn('a:ln'))
    ln.set('w', '12700')   # 1pt in EMUs
    solidFill = etree.SubElement(ln, pqn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, pqn('a:srgbClr'))
    srgbClr.set('val', '0A6EBD')

    # Caption text box below image
    txb = slide.shapes.add_textbox(x, y + IMG_H + Inches(0.04), IMG_W, CAP_H)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = caption
    run.font.size  = Pt(10)
    run.font.bold  = True
    run.font.color.rgb = BLUE_DARK
    run.font.name  = "Calibri"

    print(f"  Added: {filename}")

prs.save(str(PPTX))
print(f"  PPTX saved: {PPTX.name}\n")


# ════════════════════════════════════════════════════════════════════════════
# 2.  UPDATE DOCX — Section 11 Screenshots
# ════════════════════════════════════════════════════════════════════════════
print("Updating DOCX Section 11 (Screenshots)...")

doc = Document(str(DOCX))

# Find the paragraph that contains "11. Screenshots" or "11.1"
target_idx = None
for i, para in enumerate(doc.paragraphs):
    if "Screenshots" in para.text and ("11" in para.text or "screenshot" in para.text.lower()):
        target_idx = i
        break
    if para.text.strip().startswith("11.") and "Screenshot" in para.text:
        target_idx = i
        break

# Fallback: find any heading with "Screenshot"
if target_idx is None:
    for i, para in enumerate(doc.paragraphs):
        if "screenshot" in para.text.lower() or "11.1" in para.text:
            target_idx = i
            break

# Further fallback: insert before section 12
if target_idx is None:
    for i, para in enumerate(doc.paragraphs):
        if "12." in para.text and "Testing" in para.text:
            target_idx = i - 1
            break

print(f"  Inserting screenshots after paragraph index: {target_idx}")
print(f"  Context: '{doc.paragraphs[target_idx].text[:80] if target_idx else 'N/A'}'")

# Find the parent body element
body = doc.element.body

# Find the XML element at target_idx
if target_idx is not None:
    ref_para = doc.paragraphs[target_idx]._element
else:
    ref_para = body.findall(qn('w:p'))[-1]

def insert_after(ref, new_elem):
    ref.addnext(new_elem)

# Build a "Real Screenshots" sub-heading paragraph
def make_heading_para(text, level=3):
    p = OxmlElement('w:p')
    pPr = OxmlElement('w:pPr')
    pStyle = OxmlElement('w:pStyle')
    pStyle.set(qn('w:val'), f'Heading{level}')
    pPr.append(pStyle)
    p.append(pPr)
    r = OxmlElement('w:r')
    t = OxmlElement('w:t')
    t.text = text
    r.append(t)
    p.append(r)
    return p

def make_caption_para(text):
    p = OxmlElement('w:p')
    pPr = OxmlElement('w:pPr')
    jc = OxmlElement('w:jc'); jc.set(qn('w:val'), 'center')
    pPr.append(jc)
    p.append(pPr)
    r = OxmlElement('w:r')
    rPr = OxmlElement('w:rPr')
    b = OxmlElement('w:b'); rPr.append(b)
    color = OxmlElement('w:color'); color.set(qn('w:val'), '0A6EBD'); rPr.append(color)
    sz = OxmlElement('w:sz'); sz.set(qn('w:val'), '20'); rPr.append(sz)
    r.append(rPr)
    t = OxmlElement('w:t'); t.text = text; r.append(t)
    p.append(r)
    return p

def make_space_para():
    p = OxmlElement('w:p')
    return p

# Insert heading
heading_elem = make_heading_para("11.1  Live Application Screenshots (Deployed on AWS EC2 — ap-south-1)", level=3)
ref_para.addnext(heading_elem)
prev = heading_elem

# Insert each screenshot + caption
for filename, caption in SCREENSHOTS:
    img_path = SS_DIR / filename
    if not img_path.exists():
        print(f"  MISSING: {filename}")
        continue

    # Add paragraph with inline image using python-docx run
    img_para = doc.add_paragraph()
    img_para.alignment = 1   # CENTER
    run = img_para.add_run()
    run.add_picture(str(img_path), width=DInches(5.8))

    # Move it into position (after prev)
    img_para._element.getparent().remove(img_para._element)
    prev.addnext(img_para._element)
    prev = img_para._element

    # Caption
    cap_elem = make_caption_para(caption)
    prev.addnext(cap_elem)
    prev = cap_elem

    # Spacer
    sp = make_space_para()
    prev.addnext(sp)
    prev = sp

    print(f"  Added: {filename}")

doc.save(str(DOCX))
print(f"  DOCX saved: {DOCX.name}\n")

print("Done. Both files updated with real app screenshots.")
