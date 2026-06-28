"""
Record a live demo walkthrough of ClaimIQ at http://13.206.252.215:8501,
convert webm -> mp4, and embed the video into Slide 9 of ClaimIQ_Presentation.pptx.
"""

import pathlib, time, subprocess, shutil, sys
import imageio_ffmpeg
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE     = pathlib.Path(__file__).parent.parent
VID_DIR  = BASE / "data" / "videos"
VID_DIR.mkdir(parents=True, exist_ok=True)

APP_URL  = "http://13.206.252.215:8501"
WEBM_OUT = VID_DIR / "demo_raw.webm"
MP4_OUT  = VID_DIR / "claimiq_demo.mp4"
THUMB    = VID_DIR / "demo_thumb.png"
PPTX_PATH = BASE / "ClaimIQ_Presentation.pptx"

BLUE = RGBColor(0x0A, 0x6E, 0xBD)

# ── helpers ──────────────────────────────────────────────────────────────────
def wait(page, secs):
    time.sleep(secs)

def scroll_main(page, delta_y):
    page.evaluate(f"""
        const el = document.querySelector('.main .block-container')
                || document.querySelector('[data-testid="stMain"]')
                || document.body;
        el.scrollTop += {delta_y};
    """)
    time.sleep(1.2)

def slow_type(locator, text, delay=80):
    for ch in text:
        locator.press(ch)
        time.sleep(delay / 1000)

# ════════════════════════════════════════════════════════════════════════════
# STEP 1 — Record the walkthrough
# ════════════════════════════════════════════════════════════════════════════
print("Recording walkthrough demo...")

with sync_playwright() as p:
    browser = p.chromium.launch(headless=True, args=["--no-sandbox"])
    ctx = browser.new_context(
        viewport={"width": 1280, "height": 800},
        record_video_dir=str(VID_DIR),
        record_video_size={"width": 1280, "height": 800},
    )
    page = ctx.new_page()

    # ── 1. Load home ──────────────────────────────────────────────────────
    print("  [1/10] Loading home page...")
    page.goto(APP_URL, wait_until="networkidle", timeout=30000)
    wait(page, 5)   # show the home screen

    # ── 2. Pan sidebar ────────────────────────────────────────────────────
    print("  [2/10] Showing sidebar system status...")
    # Hover over system status items to highlight
    try:
        page.hover("text=OpenAI LLM")
    except Exception:
        pass
    wait(page, 2)
    try:
        page.hover("text=AWS S3")
    except Exception:
        pass
    wait(page, 2)
    try:
        page.hover("text=ChromaDB RAG")
    except Exception:
        pass
    wait(page, 2)

    # ── 3. Enter reviewer name ─────────────────────────────────────────────
    print("  [3/10] Entering reviewer name...")
    reviewer_typed = False
    for selector in [
        "input[placeholder='Dr. Reviewer']",
        "[data-testid='stTextInput'] input",
        "section[data-testid='stSidebar'] input",
        "input",
    ]:
        try:
            inp = page.locator(selector).first
            inp.wait_for(state="visible", timeout=3000)
            inp.click()
            wait(page, 0.5)
            slow_type(inp, "Dr. Smith")
            reviewer_typed = True
            print(f"    Reviewer typed via: {selector}")
            break
        except Exception:
            continue
    if not reviewer_typed:
        print("    WARNING: Could not type reviewer name")
    wait(page, 2)

    # ── 4. Select sample claim ────────────────────────────────────────────
    print("  [4/10] Selecting a sample claim...")
    try:
        sb = page.locator("div[data-testid='stSelectbox']").first
        sb.wait_for(state="visible", timeout=5000)
        sb.click()
        wait(page, 1.5)
        opts = page.locator("li[role='option']")
        n = opts.count()
        print(f"    {n} options found")
        if n > 1:
            opts.nth(1).click()  # first real sample
        elif n == 1:
            opts.first.click()
    except Exception as e:
        print(f"    Selectbox: {e}")
    wait(page, 2)

    # ── 5. Process Claim ──────────────────────────────────────────────────
    print("  [5/10] Clicking Process Claim...")
    try:
        btn = page.locator("button:has-text('Process Claim')").first
        btn.scroll_into_view_if_needed()
        wait(page, 1)
        btn.click()
        print("    Waiting 30s for AI pipeline...")
        wait(page, 30)
    except Exception as e:
        print(f"    Process button: {e}")
        wait(page, 5)

    # ── 6. Show claim header + OCR results ────────────────────────────────
    print("  [6/10] Scrolling through OCR results...")
    scroll_main(page, 0)    # reset scroll
    wait(page, 2)
    scroll_main(page, 300)  # claim header + OCR table visible
    wait(page, 4)

    # ── 7. Validation ─────────────────────────────────────────────────────
    print("  [7/10] Scrolling to Validation Agent results...")
    scroll_main(page, 380)
    wait(page, 4)

    # ── 8. AI Recommendation ─────────────────────────────────────────────
    print("  [8/10] Scrolling to AI Recommendation...")
    scroll_main(page, 380)
    wait(page, 4)

    # ── 9. Human Approval ─────────────────────────────────────────────────
    print("  [9/10] Scrolling to Human Approval form...")
    scroll_main(page, 400)
    wait(page, 4)
    # Try to click Approve button
    try:
        approve_btn = page.locator("button:has-text('Approve')").first
        approve_btn.scroll_into_view_if_needed()
        wait(page, 1)
        approve_btn.click()
        print("    Approved claim")
        wait(page, 5)
    except Exception as e:
        print(f"    Approve button: {e}")
        wait(page, 3)

    # ── 10. Claims History + About ────────────────────────────────────────
    print("  [10/10] Claims History and About tabs...")
    try:
        page.evaluate("window.scrollTo(0, 0)")
        wait(page, 1)
        page.locator("button[data-baseweb='tab']", has_text="Claims History").first.click()
        wait(page, 4)
    except Exception as e:
        print(f"    Claims History tab: {e}")

    try:
        page.locator("button[data-baseweb='tab']", has_text="About").first.click()
        wait(page, 4)
    except Exception as e:
        print(f"    About tab: {e}")

    # Save thumbnail before closing
    page.screenshot(path=str(THUMB))
    print("  Thumbnail saved.")

    raw_video = page.video.path()
    page.close()
    ctx.close()
    browser.close()

# Rename raw webm
if raw_video and pathlib.Path(raw_video).exists():
    shutil.move(raw_video, str(WEBM_OUT))
    print(f"  Raw video: {WEBM_OUT} ({WEBM_OUT.stat().st_size // 1024} KB)")
else:
    print("  ERROR: No video file found")
    sys.exit(1)

# ════════════════════════════════════════════════════════════════════════════
# STEP 2 — Convert webm -> mp4
# ════════════════════════════════════════════════════════════════════════════
print("\nConverting webm -> mp4...")
ffmpeg = imageio_ffmpeg.get_ffmpeg_exe()
result = subprocess.run([
    ffmpeg, "-y",
    "-i",  str(WEBM_OUT),
    "-c:v", "libx264",
    "-preset", "fast",
    "-crf", "26",          # good quality, moderate file size
    "-vf", "scale=1280:800",
    "-c:a", "aac",
    "-b:a", "128k",
    "-movflags", "+faststart",
    str(MP4_OUT)
], capture_output=True, text=True)

if result.returncode != 0:
    print("ffmpeg stderr:", result.stderr[-1000:])
    sys.exit(1)

size_mb = MP4_OUT.stat().st_size / (1024 * 1024)
print(f"  MP4 ready: {MP4_OUT.name}  ({size_mb:.1f} MB)")

# ════════════════════════════════════════════════════════════════════════════
# STEP 3 — Embed mp4 in PPTX Slide 9
# ════════════════════════════════════════════════════════════════════════════
print("\nEmbedding video into Slide 9 of PPTX...")

prs  = Presentation(str(PPTX_PATH))
slide = prs.slides[8]   # 0-based index -> slide 9

# Remove all shapes below the header (top > 1.3")
for shape in [s for s in slide.shapes if s.top > Inches(1.3)]:
    shape._element.getparent().remove(shape._element)

# ── Add the video, centred below the header ───────────────────────────────
VID_W = Inches(12.4)
VID_H = Inches(5.5)   # 1280:800 => same ratio at this width = 7.75", cap at 5.5"
VID_X = Inches((13.33 - 12.4) / 2)
VID_Y = Inches(1.35)

slide.shapes.add_movie(
    str(MP4_OUT),
    left  = VID_X,
    top   = VID_Y,
    width = VID_W,
    height= VID_H,
    poster_frame_image=str(THUMB),
    mime_type="video/mp4",
)

# Caption below video
CAP_Y = VID_Y + VID_H + Inches(0.12)
txb = slide.shapes.add_textbox(VID_X, CAP_Y, VID_W, Inches(0.35))
tf  = txb.text_frame
p   = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Live Demo — ClaimIQ deployed on AWS EC2 (ap-south-1)  |  http://13.206.252.215:8501"
run.font.size  = Pt(11)
run.font.bold  = True
run.font.color.rgb = BLUE
run.font.name  = "Calibri"

tmp = PPTX_PATH.with_name("ClaimIQ_Presentation_new.pptx")
prs.save(str(tmp))
# Try to replace the original; if locked by PowerPoint, keep the _new file
try:
    import os
    if PPTX_PATH.exists():
        os.replace(str(tmp), str(PPTX_PATH))
        final = PPTX_PATH
    else:
        final = tmp
except PermissionError:
    final = tmp
    print("  NOTE: Original PPTX is open in PowerPoint. Saved as ClaimIQ_Presentation_new.pptx")
    print("        Close PowerPoint, then rename/replace manually.")

size_pptx = final.stat().st_size / (1024 * 1024)
print(f"  PPTX saved: {final.name}  ({size_pptx:.1f} MB)")
print("\nDone. Open Slide 9 in PowerPoint to play the demo video (click the play button).")
