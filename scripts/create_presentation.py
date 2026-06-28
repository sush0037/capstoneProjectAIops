"""Generate ClaimIQ final presentation PPT from project report content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import copy
from lxml import etree

# ── Brand colours ──────────────────────────────────────────────────────────
BLUE_DARK  = RGBColor(0x0A, 0x6E, 0xBD)   # #0A6EBD  primary
BLUE_MID   = RGBColor(0x1A, 0x8C, 0xD8)   # accent
BLUE_LIGHT = RGBColor(0xE8, 0xF4, 0xFD)   # panel background
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1A, 0x1A, 0x2E)
GREY       = RGBColor(0x55, 0x55, 0x66)
GREEN      = RGBColor(0x00, 0x9B, 0x55)
ORANGE     = RGBColor(0xE8, 0x7C, 0x1E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank = prs.slide_layouts[6]   # completely blank layout


# ── helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, border_rgb=None, border_pt=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if border_rgb and border_pt:
        shape.line.color.rgb = border_rgb
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, l, t, w, h,
                 font_size=18, bold=False, italic=False,
                 color=DARK, align=PP_ALIGN.LEFT,
                 word_wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = word_wrap
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def add_para(tf, text, font_size=14, bold=False, italic=False,
             color=DARK, align=PP_ALIGN.LEFT, space_before=6,
             bullet=False, indent_level=0, font_name="Calibri"):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.level = indent_level
    run = p.add_run()
    run.text = ("• " if bullet else "") + text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return p


def slide_header(slide, title, subtitle=None):
    """Blue header bar at the top of a content slide."""
    add_rect(slide, 0, 0, 13.33, 1.15, fill_rgb=BLUE_DARK)
    add_text_box(slide, title, 0.3, 0.1, 11.0, 0.7,
                 font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.3, 0.75, 11.0, 0.35,
                     font_size=14, italic=True, color=BLUE_LIGHT, align=PP_ALIGN.LEFT)
    # thin accent line
    add_rect(slide, 0, 1.15, 13.33, 0.04, fill_rgb=BLUE_MID)
    return 1.25  # y position where content starts


def add_card(slide, l, t, w, h, title, lines, title_color=BLUE_DARK,
             bg=BLUE_LIGHT, border=BLUE_MID, font_size=12, title_size=13):
    add_rect(slide, l, t, w, h, fill_rgb=bg, border_rgb=border, border_pt=1.2)
    # card title
    add_text_box(slide, title, l+0.1, t+0.08, w-0.2, 0.3,
                 font_size=title_size, bold=True, color=title_color)
    txb = slide.shapes.add_textbox(
        Inches(l+0.1), Inches(t+0.42), Inches(w-0.2), Inches(h-0.55))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = DARK
        run.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Project Title & Team Details
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)

# Full-slide gradient-style background
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=DARK)
add_rect(s, 0, 0, 13.33, 4.2, fill_rgb=BLUE_DARK)

# Decorative accent strip
add_rect(s, 0, 4.2, 13.33, 0.08, fill_rgb=BLUE_MID)

# Title
add_text_box(s, "ClaimIQ", 0.5, 0.45, 12, 1.1,
             font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, "AI-Powered Healthcare Claim Review Assistant",
             0.5, 1.45, 12, 0.6, font_size=22, bold=False, color=BLUE_LIGHT,
             align=PP_ALIGN.CENTER)
add_text_box(s, "Intelligent · Transparent · Human-Supervised",
             0.5, 2.0, 12, 0.5, font_size=16, italic=True, color=WHITE,
             align=PP_ALIGN.CENTER)

# Submission meta
add_text_box(s, "AWS AI & ML Capstone Project  |  June 28, 2026",
             0.5, 2.75, 12, 0.4, font_size=14, color=BLUE_LIGHT,
             align=PP_ALIGN.CENTER)

# Team table header
add_rect(s, 0.7, 3.5, 11.93, 0.42, fill_rgb=BLUE_MID)
for label, xpos, wid in [
    ("#", 0.75, 0.3), ("Name", 1.35, 2.2), ("Role", 4.25, 9.35)]:
    add_text_box(s, label, xpos, 3.52, wid, 0.38,
                 font_size=13, bold=True, color=WHITE)

team = [
    ("1", "Sushilkumar Nikam", "Team Lead, Architect & AI/ML Engineer"),
    ("2", "Anupam",            "RAG Pipeline Engineer & Data Engineer"),
    ("3", "Ashish",            "AI/ML Engineer & LangGraph Workflow Architect"),
    ("4", "Bishwajit",         "Full Stack Developer, AI/ML Engineer & AWS Infrastructure Engineer"),
    ("5", "Arvind",            "AI/ML Engineer & Backend Developer"),
    ("6", "Azim",              "AI/ML Engineer & Backend Developer"),
]
row_colors = [WHITE, BLUE_LIGHT]
for i, (num, name, role) in enumerate(team):
    y = 3.92 + i * 0.52
    add_rect(s, 0.7, y, 11.93, 0.50, fill_rgb=row_colors[i % 2])
    add_text_box(s, num,  0.75, y+0.05, 0.55, 0.42, font_size=12, color=DARK)
    add_text_box(s, name, 1.35, y+0.05, 2.8,  0.42, font_size=12, bold=True, color=DARK)
    add_text_box(s, role, 4.25, y+0.05, 9.0,  0.42, font_size=11, color=DARK)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0xF7, 0xF9, 0xFC))
cy = slide_header(s, "Problem Statement", "Healthcare claim processing is broken — and expensive")

# Intro line
add_text_box(s,
    "Insurance carriers spend $30–$50 per claim in manual processing costs. "
    "Even a 10 % reduction translates to billions in annual savings.",
    0.3, cy, 12.7, 0.55, font_size=14, italic=True, color=GREY)

cy += 0.65

pain_points = [
    ("⏱  Processing Delays",    "15–30 days average per claim due to manual review queues"),
    ("❌  High Error Rates",     "Manual data entry mistakes in ICD-10 / CPT coding cause incorrect rejections"),
    ("⚖  Inconsistency",        "Different adjusters reach different decisions on identical claims"),
    ("📋  Poor Audit Trail",     "Decisions tracked via paper notes and email chains — not auditable"),
    ("🔍  Fraud Exposure",       "Manual reviewers cannot cross-reference all historical claims simultaneously"),
    ("😓  Reviewer Burnout",     "High-volume, repetitive review tasks reduce accuracy over time"),
]

cols = 3
card_w = 4.2
card_h = 1.45
gap = 0.16
start_x = 0.3
for idx, (title, desc) in enumerate(pain_points):
    col = idx % cols
    row = idx // cols
    x = start_x + col * (card_w + gap)
    y = cy + row * (card_h + 0.15)
    add_card(s, x, y, card_w, card_h, title, [desc],
             title_color=BLUE_DARK, bg=WHITE,
             border=BLUE_MID, font_size=12, title_size=12)

# Core problem callout at bottom
add_rect(s, 0.3, 6.3, 12.73, 0.92, fill_rgb=BLUE_DARK)
add_text_box(s,
    "Core Gap: No unified system that can extract, validate, retrieve policy context, "
    "recommend a decision, and present it to a human — in a single automated workflow.",
    0.5, 6.36, 12.33, 0.8, font_size=13, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Selected Theme and Domain
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0xF7, 0xF9, 0xFC))
cy = slide_header(s, "Selected Theme & Domain", "")

# Two big cards side by side
for x, bg, bdr, title, icon, lines in [
    (0.3,  WHITE,      BLUE_MID, "Domain",
     "🏥  Healthcare & Insurance Technology (HealthTech / InsurTech)",
     ["Insurance claim review automation",
      "HIPAA-compliant data processing",
      "AI-assisted decision support for claims adjusters",
      "Fraud detection and compliance enforcement"]),
    (6.75, BLUE_LIGHT, BLUE_DARK, "Theme",
     "🤖  Agentic AI with Human-in-the-Loop Oversight",
     ["Multiple specialised AI agents collaborate in a defined workflow",
      "Each agent performs a distinct role: extract → validate → recommend",
      "LangGraph orchestrates the multi-agent pipeline as a StateGraph",
      "Mandatory human approval gate before every final decision",
      "All decisions are auditable and overridable by human reviewers"]),
]:
    add_rect(s, x, cy, 6.2, 2.8, fill_rgb=bg, border_rgb=bdr, border_pt=1.5)
    add_text_box(s, title, x+0.15, cy+0.1, 5.9, 0.3,
                 font_size=11, bold=True, color=GREY)
    add_text_box(s, icon,  x+0.15, cy+0.35, 5.9, 0.5,
                 font_size=14, bold=True, color=BLUE_DARK)
    txb = slide.shapes.add_textbox if False else \
          s.shapes.add_textbox(Inches(x+0.15), Inches(cy+0.9),
                               Inches(5.8), Inches(1.8))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True; first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_before = Pt(4)
        run = p.add_run(); run.text = "• " + line
        run.font.size = Pt(12); run.font.color.rgb = DARK
        run.font.name = "Calibri"

# Theme alignment box
add_rect(s, 0.3, cy+3.0, 12.73, 1.7, fill_rgb=WHITE, border_rgb=BLUE_MID, border_pt=1)
add_text_box(s, "Why This Domain + Theme Combination?",
             0.5, cy+3.08, 12.33, 0.35, font_size=13, bold=True, color=BLUE_DARK)
bullets = [
    "Healthcare claims are document-heavy, rule-intensive, and high-stakes — ideal for structured multi-agent AI.",
    "Policy rules change frequently — RAG lets the system stay current without retraining the model.",
    "Financial and patient-care consequences demand human oversight — the Human-in-the-Loop gate ensures AI is advisory, not authoritative.",
]
txb2 = s.shapes.add_textbox(Inches(0.5), Inches(cy+3.48), Inches(12.33), Inches(1.15))
txb2.word_wrap = True; tf2 = txb2.text_frame; tf2.word_wrap = True; first = True
for b in bullets:
    p = tf2.paragraphs[0] if first else tf2.add_paragraph(); first = False
    p.space_before = Pt(3)
    run = p.add_run(); run.text = "• " + b
    run.font.size = Pt(12); run.font.color.rgb = DARK; run.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Solution Overview
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0xF7, 0xF9, 0xFC))
cy = slide_header(s, "Solution Overview", "ClaimIQ — end-to-end AI claim review in under 20 seconds")

add_text_box(s,
    "ClaimIQ is an AI-powered, multi-agent healthcare claim review system that orchestrates "
    "four specialised AI agents using LangGraph, enhanced by a RAG knowledge base of insurance "
    "policy documents, culminating in a human-supervised final decision.",
    0.3, cy, 12.73, 0.65, font_size=13, italic=True, color=GREY)

cy += 0.75

# 4-agent pipeline flow
agents = [
    ("1", "OCR Agent",            "pdfplumber +\nGPT-4o-mini",   "Extracts 15 structured\nfields from PDF"),
    ("2", "Validation Agent",     "Rule-based\nlogic",           "Checks 12 fields,\ncodes & dates"),
    ("3", "Recommendation\nAgent","RAG + GPT-4o-mini",          "Policy-grounded\nAPPROVE/REJECT/REVIEW"),
    ("4", "Human Approval\nGate", "Streamlit UI",               "Reviewer confirms\nor overrides AI"),
]
box_w = 2.9; gap_x = 0.25; start_x = 0.3
arrow_y = cy + 0.9
for i, (num, name, tech, desc) in enumerate(agents):
    x = start_x + i * (box_w + gap_x)
    # main box
    add_rect(s, x, cy, box_w, 1.85, fill_rgb=BLUE_DARK if i < 3 else GREEN)
    add_text_box(s, f"Agent {num}" if i < 3 else "Human Gate",
                 x+0.1, cy+0.08, box_w-0.2, 0.3,
                 font_size=10, bold=False, color=BLUE_LIGHT if i < 3 else WHITE)
    add_text_box(s, name, x+0.1, cy+0.32, box_w-0.2, 0.55,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, tech, x+0.1, cy+0.88, box_w-0.2, 0.45,
                 font_size=10, italic=True, color=BLUE_LIGHT if i < 3 else WHITE,
                 align=PP_ALIGN.CENTER)
    add_text_box(s, desc, x+0.1, cy+1.38, box_w-0.2, 0.45,
                 font_size=10, color=WHITE, align=PP_ALIGN.CENTER)
    # arrow between boxes
    if i < 3:
        ax = x + box_w + 0.03
        add_text_box(s, "▶", ax, cy+0.75, 0.22, 0.4,
                     font_size=16, bold=True, color=BLUE_MID, align=PP_ALIGN.CENTER)

cy += 2.05

# Design principles row
principles = [
    ("🔍 Transparent", "Every recommendation cites the exact policy clause"),
    ("👤 Human Control", "AI never finalises a claim unilaterally"),
    ("🛡 Graceful Degradation", "Works even if OpenAI or S3 is unavailable"),
    ("📁 Full Audit Trail", "Every decision stored as timestamped JSON"),
]
card_w2 = 3.1; gap2 = 0.14
for i, (title, desc) in enumerate(principles):
    x = 0.3 + i * (card_w2 + gap2)
    add_card(s, x, cy, card_w2, 1.55, title, [desc],
             title_color=BLUE_DARK, bg=WHITE, border=BLUE_MID,
             font_size=12, title_size=12)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Architecture Diagram
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0xF7, 0xF9, 0xFC))
cy = slide_header(s, "System Architecture", "Three-layer design: UI → Orchestration → Services")

# Layer labels on left
for y_pos, label, clr in [
    (1.45, "USER LAYER",          BLUE_DARK),
    (3.05, "ORCHESTRATION LAYER", BLUE_DARK),
    (5.0,  "SERVICES LAYER",      BLUE_DARK),
]:
    add_rect(s, 0.15, y_pos, 1.05, 1.35, fill_rgb=BLUE_DARK)
    txb = s.shapes.add_textbox(Inches(0.18), Inches(y_pos+0.1),
                               Inches(0.99), Inches(1.15))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = label
    run.font.size = Pt(9); run.font.bold = True
    run.font.color.rgb = WHITE; run.font.name = "Calibri"

# ─── User Layer ───────────────────────────────────────────────────────────
add_rect(s, 1.35, 1.4, 11.65, 1.45, fill_rgb=WHITE, border_rgb=BLUE_MID, border_pt=1)
for i, (tab, desc) in enumerate([
    ("Process New Claim", "Upload PDF / select sample"),
    ("Review & Approve",  "Human Approval Gate"),
    ("Claims History",    "Dashboard & metrics"),
    ("About / Info",      "System status"),
]):
    x = 1.5 + i * 2.85
    add_rect(s, x, 1.52, 2.6, 0.88, fill_rgb=BLUE_LIGHT, border_rgb=BLUE_MID, border_pt=0.8)
    add_text_box(s, tab,  x+0.1, 1.56, 2.4, 0.38, font_size=11, bold=True, color=BLUE_DARK)
    add_text_box(s, desc, x+0.1, 1.9,  2.4, 0.38, font_size=10, color=GREY)

# down arrow
add_text_box(s, "▼", 6.5, 2.88, 0.4, 0.3, font_size=14, bold=True,
             color=BLUE_MID, align=PP_ALIGN.CENTER)

# ─── Orchestration Layer ──────────────────────────────────────────────────
add_rect(s, 1.35, 3.0, 11.65, 1.5, fill_rgb=WHITE, border_rgb=BLUE_MID, border_pt=1)
add_text_box(s, "LangGraph StateGraph  (ClaimState)", 1.5, 3.05, 5.0, 0.3,
             font_size=11, bold=True, color=BLUE_DARK)
for i, (name, clr) in enumerate([
    ("OCR Agent",            BLUE_DARK),
    ("Validation Agent",     BLUE_DARK),
    ("Recommendation Agent", BLUE_DARK),
    ("Human Gate",           GREEN),
]):
    x = 1.5 + i * 2.85
    add_rect(s, x, 3.38, 2.6, 0.88, fill_rgb=clr)
    add_text_box(s, name, x+0.08, 3.55, 2.44, 0.55,
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text_box(s, "▶", x+2.6+0.03, 3.72, 0.22, 0.3,
                     font_size=13, bold=True, color=BLUE_MID, align=PP_ALIGN.CENTER)

# down arrows
for x_a in [2.7, 6.5, 10.2]:
    add_text_box(s, "▼", x_a, 4.5, 0.4, 0.3, font_size=14, bold=True,
                 color=BLUE_MID, align=PP_ALIGN.CENTER)

# ─── Services Layer ───────────────────────────────────────────────────────
svc_items = [
    (1.5,  BLUE_LIGHT, BLUE_MID, "AI Services",
     ["OpenAI GPT-4o-mini", "text-embedding-3-small"]),
    (5.25, BLUE_LIGHT, BLUE_MID, "Knowledge Base",
     ["ChromaDB Vector Store", "3 Policy PDFs (~150 chunks)"]),
    (9.0,  BLUE_LIGHT, BLUE_MID, "Storage",
     ["Local JSON (data/outputs/)", "AWS S3 (optional cloud)"]),
]
for x, bg, bdr, title, lines in svc_items:
    add_rect(s, x, 4.82, 3.5, 1.5, fill_rgb=bg, border_rgb=bdr, border_pt=1)
    add_text_box(s, title, x+0.12, 4.88, 3.26, 0.32,
                 font_size=12, bold=True, color=BLUE_DARK)
    for j, line in enumerate(lines):
        add_text_box(s, "• " + line, x+0.12, 5.22+j*0.38, 3.26, 0.36,
                     font_size=11, color=DARK)

# Data flow caption
add_rect(s, 0.15, 6.44, 12.9, 0.85, fill_rgb=BLUE_DARK)
add_text_box(s,
    "Data Flow: PDF Upload → OCR → Validation → RAG Retrieval → LLM Recommendation → Human Approval → S3 Storage",
    0.3, 6.5, 12.7, 0.73, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Agent Workflow
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0xF7, 0xF9, 0xFC))
cy = slide_header(s, "Agent Workflow Design", "LangGraph StateGraph — four nodes, one shared ClaimState")

agents_detail = [
    (BLUE_DARK, "Agent 1 — OCR Agent", "src/agents/ocr_agent.py", [
        "Opens PDF with pdfplumber, extracts raw text",
        "Calls GPT-4o-mini with response_format=json_object",
        "Returns 15 structured fields (patient, provider, codes, amounts)",
        "Sets confidence: HIGH (12+ fields) / MEDIUM / LOW",
    ]),
    (BLUE_MID,  "Agent 2 — Validation Agent", "src/agents/validation_agent.py", [
        "Checks 6 mandatory fields — any missing → severity HIGH",
        "Checks 6 important fields — missing → severity MEDIUM",
        "Validates ICD-10, CPT, NPI formats via regex",
        "Validates date logic: treatment ≤ claim date, filed within 365 days",
    ]),
    (BLUE_DARK, "Agent 3 — Recommendation Agent", "src/agents/recommendation_agent.py", [
        "Fast-path rules: duplicate, dental, experimental, >$25K → instant decision",
        "Retrieves top-5 policy chunks from ChromaDB via cosine similarity",
        "Calls GPT-4o-mini with claim data + policy context",
        "Returns APPROVE / REJECT / MANUAL_REVIEW with citations",
    ]),
    (GREEN,     "Human Approval Gate", "app.py + finalize_claim()", [
        "Pipeline pauses at awaiting_human_approval = True",
        "Reviewer sees extracted fields, validation results, AI recommendation",
        "Can Approve, Reject, or Override (override requires written justification)",
        "finalize_claim() persists decision as timestamped JSON + S3 upload",
    ]),
]

card_h = 1.42; card_w = 12.7; start_y = cy + 0.05
for i, (clr, title, file, bullets) in enumerate(agents_detail):
    y = start_y + i * (card_h + 0.1)
    add_rect(s, 0.3, y, 0.18, card_h, fill_rgb=clr)   # colour tab
    add_rect(s, 0.52, y, card_w-0.22, card_h, fill_rgb=WHITE, border_rgb=BLUE_MID, border_pt=0.8)
    add_text_box(s, title, 0.65, y+0.06, 6.5, 0.32,
                 font_size=13, bold=True, color=clr)
    add_text_box(s, file,  7.5,  y+0.06, 5.5, 0.32,
                 font_size=10, italic=True, color=GREY, align=PP_ALIGN.RIGHT)
    txb = s.shapes.add_textbox(Inches(0.65), Inches(y+0.42), Inches(12.3), Inches(0.95))
    txb.word_wrap = True; tf = txb.text_frame; tf.word_wrap = True; first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_before = Pt(2)
        run = p.add_run(); run.text = "• " + b
        run.font.size = Pt(11); run.font.color.rgb = DARK; run.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — RAG Pipeline
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0xF7, 0xF9, 0xFC))
cy = slide_header(s, "RAG Pipeline", "Retrieval-Augmented Generation — grounding AI in current policy documents")

# Why RAG
add_text_box(s,
    "Policy documents change — RAG lets the system use current rules without retraining the model. "
    "Every recommendation includes the exact policy clause that justified it.",
    0.3, cy, 12.73, 0.52, font_size=13, italic=True, color=GREY)
cy += 0.62

# ── Ingestion pipeline ────────────────────────────────────────────────────
add_text_box(s, "INGESTION PIPELINE  (scripts/ingest_documents.py)",
             0.3, cy, 12.73, 0.32, font_size=12, bold=True, color=BLUE_DARK)
cy += 0.38

ingest_steps = [
    ("📄 Policy PDFs\n(3 documents)", BLUE_LIGHT),
    ("📝 Text Extraction\n(pdfplumber)", WHITE),
    ("✂️  Text Chunking\n800 chars / 100 overlap", WHITE),
    ("🔢 OpenAI Embeddings\ntext-embedding-3-small", WHITE),
    ("🗄 ChromaDB\n~150 chunks persisted", BLUE_LIGHT),
]
bw = 2.3; bh = 1.0; sx = 0.3
for i, (label, bg) in enumerate(ingest_steps):
    x = sx + i * (bw + 0.18)
    add_rect(s, x, cy, bw, bh, fill_rgb=bg, border_rgb=BLUE_MID, border_pt=1)
    add_text_box(s, label, x+0.08, cy+0.1, bw-0.16, bh-0.15,
                 font_size=11, bold=False, color=DARK, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text_box(s, "→", x+bw+0.02, cy+0.3, 0.16, 0.4,
                     font_size=13, bold=True, color=BLUE_MID, align=PP_ALIGN.CENTER)
cy += 1.15

# ── Retrieval pipeline ────────────────────────────────────────────────────
add_text_box(s, "RETRIEVAL PIPELINE  (at claim review time)",
             0.3, cy, 12.73, 0.32, font_size=12, bold=True, color=BLUE_DARK)
cy += 0.38

retrieval_steps = [
    ("📋 Claim Fields\n(extracted)", BLUE_LIGHT),
    ("🔍 Build Query String\ndiagnosis + procedure\n+ amount", WHITE),
    ("🔢 Query Embedding\ntext-embedding-3-small", WHITE),
    ("📊 ChromaDB Search\ntop-5 cosine similarity", WHITE),
    ("💡 Inject into\nLLM Prompt", BLUE_LIGHT),
]
for i, (label, bg) in enumerate(retrieval_steps):
    x = sx + i * (bw + 0.18)
    add_rect(s, x, cy, bw, 1.1, fill_rgb=bg, border_rgb=BLUE_MID, border_pt=1)
    add_text_box(s, label, x+0.08, cy+0.08, bw-0.16, 0.95,
                 font_size=11, color=DARK, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text_box(s, "→", x+bw+0.02, cy+0.35, 0.16, 0.4,
                     font_size=13, bold=True, color=BLUE_MID, align=PP_ALIGN.CENTER)
cy += 1.25

# Embedding model table
add_rect(s, 0.3, cy, 5.8, 1.48, fill_rgb=WHITE, border_rgb=BLUE_MID, border_pt=1)
add_text_box(s, "Embedding Model Details", 0.45, cy+0.06, 5.5, 0.3,
             font_size=12, bold=True, color=BLUE_DARK)
for j, (k, v) in enumerate([
    ("Model", "text-embedding-3-small"),
    ("Dimensions", "1,536"),
    ("Retrieval top-k", "5  (cosine similarity)"),
    ("Chunk size", "800 chars / 100 overlap"),
]):
    add_text_box(s, f"{k}:", 0.45, cy+0.42+j*0.25, 1.7,  0.24, font_size=11, bold=True, color=DARK)
    add_text_box(s, v,      2.2,  cy+0.42+j*0.25, 3.8,  0.24, font_size=11, color=DARK)

# Fallback strategy box
add_rect(s, 6.6, cy, 6.5, 1.48, fill_rgb=BLUE_LIGHT, border_rgb=BLUE_MID, border_pt=1)
add_text_box(s, "🛡 Fallback Strategy", 6.75, cy+0.06, 6.2, 0.3,
             font_size=12, bold=True, color=BLUE_DARK)
add_text_box(s,
    "If ChromaDB is unavailable or index is empty, the retriever returns "
    "a hard-coded FALLBACK_POLICY_CONTEXT with key rules extracted from "
    "all three policy documents — ensuring the Recommendation Agent always "
    "has some policy grounding.",
    6.75, cy+0.42, 6.2, 1.0, font_size=11, color=DARK)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — AWS Deployment Overview
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0xF7, 0xF9, 0xFC))
cy = slide_header(s, "AWS Deployment Overview", "Two patterns: EC2 Docker Compose (dev) and ECS Fargate (production)")

# Pattern A
add_rect(s, 0.3, cy, 6.1, 4.85, fill_rgb=WHITE, border_rgb=BLUE_DARK, border_pt=1.5)
add_rect(s, 0.3, cy, 6.1, 0.38, fill_rgb=BLUE_DARK)
add_text_box(s, "Pattern A — EC2 + Docker Compose", 0.45, cy+0.04, 5.9, 0.3,
             font_size=12, bold=True, color=WHITE)
for j, (label, val) in enumerate([
    ("Instance",  "Ubuntu 22.04, t3.medium (2 vCPU / 4 GB)"),
    ("Storage",   "20 GB gp3 EBS"),
    ("Region",    "us-east-1"),
    ("Access",    "Port 8501 (Streamlit) / 80+443 (Nginx)"),
    ("Use Case",  "Development, demos, low cost (~$30/mo)"),
]):
    y = cy + 0.48 + j * 0.4
    add_text_box(s, label+":", 0.45, y, 1.4, 0.35, font_size=11, bold=True, color=BLUE_DARK)
    add_text_box(s, val,       1.9,  y, 4.4, 0.35, font_size=11, color=DARK)

# Docker compose snippet
add_rect(s, 0.45, cy+2.55, 5.8, 1.6, fill_rgb=RGBColor(0xF0, 0xF4, 0xF8))
add_text_box(s,
    "services:\n  claimiq:\n    build: .\n    ports: [\"8501:8501\"]\n"
    "    env_file: .env\n    volumes: [./data:/app/data]\n    restart: unless-stopped",
    0.55, cy+2.6, 5.6, 1.5, font_size=9, color=DARK, font_name="Courier New")

# Pattern B
add_rect(s, 6.93, cy, 6.1, 4.85, fill_rgb=WHITE, border_rgb=GREEN, border_pt=1.5)
add_rect(s, 6.93, cy, 6.1, 0.38, fill_rgb=GREEN)
add_text_box(s, "Pattern B — ECS Fargate + ALB", 7.08, cy+0.04, 5.9, 0.3,
             font_size=12, bold=True, color=WHITE)
for j, (label, val) in enumerate([
    ("Compute",   "ECS Fargate — 1 vCPU / 2 GB Memory"),
    ("Registry",  "Amazon ECR — Docker image store"),
    ("Secrets",   "AWS Secrets Manager (OpenAI key)"),
    ("Load Bal.", "ALB → port 80/443 → target 8501"),
    ("Logs",      "CloudWatch /ecs/claimiq"),
    ("Use Case",  "Production auto-scaling (~$40/mo)"),
]):
    y = cy + 0.48 + j * 0.4
    add_text_box(s, label+":", 7.08, y, 1.45, 0.35, font_size=11, bold=True, color=GREEN)
    add_text_box(s, val,       8.58, y, 4.35, 0.35, font_size=11, color=DARK)

# AWS services grid
add_text_box(s, "AWS Services Used", 0.3, cy+5.0, 12.73, 0.32,
             font_size=12, bold=True, color=BLUE_DARK)
svcs = ["EC2", "ECS Fargate", "ECR", "S3", "Secrets Manager", "IAM", "ALB", "CloudWatch"]
for i, svc in enumerate(svcs):
    x = 0.3 + i * 1.62
    add_rect(s, x, cy+5.38, 1.5, 0.52, fill_rgb=BLUE_LIGHT, border_rgb=BLUE_MID, border_pt=0.8)
    add_text_box(s, svc, x+0.06, cy+5.44, 1.38, 0.4,
                 font_size=10, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Demo Screenshots
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0xF7, 0xF9, 0xFC))
cy = slide_header(s, "Demo Screenshots", "Streamlit web application — five key screens")

screens = [
    ("1. Home — Process New Claim",
     ["File upload drag-and-drop zone",
      "OR select from 10 pre-built sample claims",
      "Sidebar: OpenAI ✅  ChromaDB ✅  S3 ✅",
      "Quick Actions: Generate Data / Build RAG Index"]),
    ("2. AI Processing Pipeline",
     ["4-step stepper: OCR → Validate → Recommend → Human",
      "Extracted fields card with ✅ / ⚠️ per field",
      "Validation severity badge (LOW / MEDIUM / HIGH)",
      "AI recommendation with confidence + citations"]),
    ("3. Human Approval Interface",
     ["Radio: Approve | Reject | Override Approve | Override Reject",
      "Notes textarea (mandatory for overrides)",
      "Submit button disabled until selection made",
      "Override warning banner with justification prompt"]),
    ("4. Claims History Dashboard",
     ["Total / Approved / Rejected / Manual Review counts",
      "AI Agreement Rate metric",
      "Sortable colour-coded table (green/red/yellow)",
      "Export functionality for audit records"]),
    ("5. About / System Info",
     ["Technology stack with version info",
      "Pipeline diagram",
      "Supported claim scenarios and expected outcomes",
      "Team credits"]),
]

card_w3 = 2.45; card_h3 = 2.55; gap3 = 0.12; sx3 = 0.3
for i, (title, bullets) in enumerate(screens):
    col = i % 3; row = i // 3
    x = sx3 + col * (card_w3 + gap3)
    y = cy + 0.05 + row * (card_h3 + 0.15)
    # mock screen area
    add_rect(s, x, y, card_w3, 1.05, fill_rgb=DARK)
    add_rect(s, x+0.06, y+0.06, card_w3-0.12, 0.18, fill_rgb=BLUE_DARK)
    for dot_x in [x+0.09, x+0.19, x+0.29]:
        add_rect(s, dot_x, y+0.09, 0.07, 0.07, fill_rgb=WHITE)
    # fake content lines
    for ln in range(4):
        lw = card_w3 - 0.2 - (ln % 2) * 0.4
        add_rect(s, x+0.1, y+0.32+ln*0.17, lw, 0.1,
                 fill_rgb=BLUE_LIGHT if ln < 2 else RGBColor(0xCC, 0xCC, 0xCC))
    # caption
    add_text_box(s, title, x, y+1.1, card_w3, 0.32,
                 font_size=10, bold=True, color=BLUE_DARK)
    txb = s.shapes.add_textbox(Inches(x), Inches(y+1.45), Inches(card_w3), Inches(1.05))
    txb.word_wrap = True; tf = txb.text_frame; tf.word_wrap = True; first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_before = Pt(1)
        run = p.add_run(); run.text = "• " + b
        run.font.size = Pt(9); run.font.color.rgb = DARK; run.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Key Outcomes
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0xF7, 0xF9, 0xFC))
cy = slide_header(s, "Key Outcomes & Results", "100 % accuracy on all 10 synthetic claim scenarios")

# Big stat boxes
stats = [
    ("94.1 %", "Average OCR\nField Extraction Rate"),
    ("100 %",  "Validation Severity\nAccuracy (10/10)"),
    ("100 %",  "Recommendation\nAccuracy (10/10)"),
    ("< 20 s", "End-to-End\nProcessing Time"),
]
bw4 = 2.95; bh4 = 1.45
for i, (num, label) in enumerate(stats):
    x = 0.3 + i * (bw4 + 0.2)
    add_rect(s, x, cy, bw4, bh4, fill_rgb=BLUE_DARK)
    add_text_box(s, num,   x+0.1, cy+0.15, bw4-0.2, 0.65,
                 font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, label, x+0.1, cy+0.82, bw4-0.2, 0.55,
                 font_size=11, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
cy += 1.6

# Business impact comparison table
add_text_box(s, "Business Impact — Manual Process vs ClaimIQ",
             0.3, cy, 12.73, 0.32, font_size=13, bold=True, color=BLUE_DARK)
cy += 0.38

headers = ["Step", "Manual Process", "ClaimIQ (AI)", "Improvement"]
col_widths = [3.4, 3.1, 2.9, 3.2]
col_x = [0.3, 3.75, 6.9, 9.83]
add_rect(s, 0.3, cy, 12.73, 0.38, fill_rgb=BLUE_DARK)
for hdr, x, w in zip(headers, col_x, col_widths):
    add_text_box(s, hdr, x+0.08, cy+0.05, w-0.12, 0.28,
                 font_size=11, bold=True, color=WHITE)
cy += 0.38

rows = [
    ("Extract claim fields",      "5–10 minutes",   "8–12 seconds",  "~50× faster"),
    ("Validate rules",            "3–5 minutes",    "< 1 second",    "300× faster"),
    ("Look up policy context",    "5–15 minutes",   "2–4 seconds",   "~200× faster"),
    ("Total before human review", "15–30 minutes",  "~15–20 seconds","~100× faster"),
    ("Audit trail completeness",  "Inconsistent",   "100 % JSON",    "Fully auditable"),
    ("Decision consistency",      "Variable",       "Deterministic", "Fully reproducible"),
]
row_colors2 = [WHITE, BLUE_LIGHT]
for ri, row in enumerate(rows):
    rc = row_colors2[ri % 2]
    add_rect(s, 0.3, cy, 12.73, 0.42, fill_rgb=rc)
    for ci, (val, x, w) in enumerate(zip(row, col_x, col_widths)):
        clr = GREEN if ci == 3 else DARK
        bold = ci == 3
        add_text_box(s, val, x+0.08, cy+0.06, w-0.12, 0.3,
                     font_size=10, bold=bold, color=clr)
    cy += 0.42


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Challenges & Learnings
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0xF7, 0xF9, 0xFC))
cy = slide_header(s, "Challenges & Learnings", "Seven production-grade challenges resolved during development")

challenges = [
    ("Image-based PDFs",         "pdfplumber extracts no text from scanned PDFs",
     "Text-length guard + graceful MANUAL_REVIEW fallback"),
    ("LangGraph Serialization",  "Non-primitive types in ClaimState crashed the pipeline",
     "Agents store only strings/dicts/lists; json.dumps validation at each node"),
    ("ChromaDB Corruption",      "Interrupted ingestion left partial, misleading index",
     "--reset flag to wipe and recreate; chunk count validation post-ingestion"),
    ("LLM JSON Parsing",         "GPT-4o-mini wrapped JSON in markdown code fences",
     "3-step cleaning: direct parse → strip fences → regex extract → fallback"),
    ("Session State Race Cond.", "Two concurrent submissions interleaved UI results",
     "processing_lock flag; claim_id correlation; explicit state reset on new claim"),
    ("EC2 Docker Permissions",   "ubuntu user not in docker group on fresh EC2",
     "usermod -aG docker ubuntu step added to DEPLOYMENT.md checklist"),
    ("Cold Start Latency",       "First claim took 25–40 s — blank spinner confused users",
     "Module-level initialisation; Streamlit progress spinner with status text"),
]

cw = 4.1; ch = 1.3; gap_c = 0.12
positions = [
    (0.28, cy), (4.50, cy), (8.72, cy),
    (0.28, cy+1.45), (4.50, cy+1.45), (8.72, cy+1.45),
    (0.28, cy+2.9),
]
for (x, y), (title, prob, sol) in zip(positions, challenges):
    add_rect(s, x, y, cw, ch, fill_rgb=WHITE, border_rgb=BLUE_MID, border_pt=1)
    add_rect(s, x, y, cw, 0.28, fill_rgb=BLUE_DARK)
    add_text_box(s, title, x+0.08, y+0.03, cw-0.16, 0.22,
                 font_size=10, bold=True, color=WHITE)
    add_text_box(s, "❌ " + prob, x+0.08, y+0.32, cw-0.16, 0.38,
                 font_size=9.5, color=RGBColor(0xCC, 0x33, 0x33))
    add_text_box(s, "✅ " + sol, x+0.08, y+0.72, cw-0.16, 0.52,
                 font_size=9.5, color=GREEN)

# Learnings box
ly = cy + 4.38
add_rect(s, 0.28, ly, 12.77, 1.85, fill_rgb=BLUE_LIGHT, border_rgb=BLUE_DARK, border_pt=1)
add_text_box(s, "Key Learnings", 0.42, ly+0.08, 12.5, 0.3,
             font_size=13, bold=True, color=BLUE_DARK)
learnings = [
    "LangGraph StateGraph requires all state values to be JSON-serialisable — design state contracts early",
    "RAG citation quality depends heavily on chunk size and overlap — 800/100 chars proved optimal for policy docs",
    "Graceful degradation must be designed in from day one — every external call needs a fallback path",
    "Human-in-the-loop UI needs clear confidence signals — confidence badges and policy citations drive reviewer trust",
]
txb = s.shapes.add_textbox(Inches(0.42), Inches(ly+0.44), Inches(12.5), Inches(1.35))
txb.word_wrap = True; tf = txb.text_frame; tf.word_wrap = True; first = True
for learn in learnings:
    p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
    p.space_before = Pt(3)
    run = p.add_run(); run.text = "• " + learn
    run.font.size = Pt(11); run.font.color.rgb = DARK; run.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Future Scope
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=DARK)
add_rect(s, 0, 0, 13.33, 1.15, fill_rgb=BLUE_DARK)
add_rect(s, 0, 1.15, 13.33, 0.04, fill_rgb=BLUE_MID)
add_text_box(s, "Future Scope & Enhancements", 0.3, 0.1, 12.7, 0.75,
             font_size=28, bold=True, color=WHITE)
add_text_box(s, "Building on ClaimIQ's foundation — the roadmap ahead",
             0.3, 0.78, 12.7, 0.34, font_size=14, italic=True, color=BLUE_LIGHT)

futures = [
    ("🔍  AWS Textract Integration",
     "Replace pdfplumber with AWS Textract to handle image-based (scanned) PDFs "
     "without requiring a digital text layer — critical for real-world claim submissions."),
    ("🔄  Live Policy Document Updates",
     "Trigger automatic RAG re-ingestion via S3 event notifications when new policy "
     "PDFs are uploaded — keeping the knowledge base always current with zero manual intervention."),
    ("🧠  Feedback Loop & Fine-Tuning",
     "Capture human override decisions as training signal to iteratively improve the "
     "Recommendation Agent's prompt and reduce override rates over time."),
    ("🏥  Multi-Plan Support",
     "Parameterise the RAG retriever to filter by insurance plan code — enabling a "
     "single ClaimIQ instance to serve multiple insurance products simultaneously."),
    ("📊  Explainability Dashboard",
     "Visualise which specific policy chunks influenced each recommendation and display "
     "their cosine similarity scores — full AI decision explainability for compliance teams."),
    ("🔔  Notification Integration",
     "Email / Slack alerts to reviewers when high-priority claims arrive (experimental "
     "treatments, amounts > $25K) — reducing review turnaround for critical cases."),
]

card_w5 = 4.05; card_h5 = 1.88; gap5 = 0.18; sx5 = 0.3
for i, (title, desc) in enumerate(futures):
    col = i % 3; row = i // 3
    x = sx5 + col * (card_w5 + gap5)
    y = 1.35 + row * (card_h5 + 0.18)
    add_rect(s, x, y, card_w5, card_h5,
             fill_rgb=RGBColor(0x1E, 0x2A, 0x3A),
             border_rgb=BLUE_MID, border_pt=1.2)
    add_text_box(s, title, x+0.14, y+0.12, card_w5-0.28, 0.42,
                 font_size=13, bold=True, color=BLUE_LIGHT)
    add_text_box(s, desc, x+0.14, y+0.6, card_w5-0.28, 1.2,
                 font_size=11, color=RGBColor(0xCC, 0xD6, 0xE0))

# Bottom tagline
add_rect(s, 0, 6.85, 13.33, 0.65, fill_rgb=BLUE_DARK)
add_text_box(s,
    "ClaimIQ Squad  •  Sushilkumar | Anupam | Ashish | Bishwajit | Arvind | Azim  •  June 2026",
    0, 6.9, 13.33, 0.55, font_size=12, color=WHITE, align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────────
import pathlib
out = pathlib.Path(__file__).parent.parent / "ClaimIQ_Presentation.pptx"
prs.save(str(out))
print(f"Saved: {out}")
